"""
Local Study Search Engine (TF-IDF)

Features:
- Desktop GUI for searching local study materials
- Supports docx / pptx / txt / md
- Supports pdf if pypdf or PyPDF2 is installed
- Reads DOCX paragraphs and tables
- Uses exact full-text matching + TF-IDF keyword ranking
- Can open files, copy paths, and export results
- Can also run in command line mode

Usage:
  python study_search.py
  python study_search.py gui
  python study_search.py index <folder_path> [index.json]
  python study_search.py search [index.json]
  python study_search.py query <index.json> "<your query>" [-k 5]

Recommended installation:
  pip install python-docx python-pptx pypdf
"""

from __future__ import annotations

import os
import re
import sys
import json
import math
import time
import threading
import subprocess
from dataclasses import dataclass, asdict
from datetime import datetime
from typing import Callable, Dict, List, Tuple, Optional

try:
    import tkinter as tk
    from tkinter import filedialog, messagebox, ttk
except Exception:
    tk = None
    filedialog = None
    messagebox = None
    ttk = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pypdf import PdfReader  # type: ignore
except Exception:
    try:
        from PyPDF2 import PdfReader  # type: ignore
    except Exception:
        PdfReader = None


SUPPORTED_EXT = {".txt", ".md", ".docx", ".pptx"}

if PdfReader is not None:
    SUPPORTED_EXT.add(".pdf")


STOPWORDS = {
    "a", "an", "the", "and", "or", "but", "if", "then", "else",
    "to", "of", "in", "on", "for", "with", "as", "by", "at", "from",
    "is", "are", "was", "were", "be", "been", "being",
    "it", "this", "that", "these", "those",
    "we", "you", "they", "i", "he", "she",
    "can", "could", "should", "would", "may", "might",
    "not", "no", "yes",
}

TOKEN_RE = re.compile(r"[a-z0-9]+", re.IGNORECASE)
SEARCH_TERM_RE = re.compile(r"[a-z0-9]+|[\u4e00-\u9fff]+", re.IGNORECASE)


@dataclass
class Document:
    doc_id: str
    path: str
    title: str
    text: str
    n_chars: int
    n_tokens: int


@dataclass
class IndexMeta:
    created_at: float
    root_folder: str
    n_docs: int
    note: str = "TF-IDF index for local study materials"


@dataclass
class TfIdfIndex:
    meta: IndexMeta
    documents: List[Document]
    idf: Dict[str, float]
    vectors: Dict[str, Dict[str, float]]
    norms: Dict[str, float]


@dataclass
class SearchResult:
    score: float
    document: Document
    snippet: str
    matches: int
    source: str


# ---------------------------
# Text extraction
# ---------------------------

def read_text_file(path: str) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception:
            continue
    return ""


def extract_docx(path: str) -> str:
    if docx is None:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx")

    d = docx.Document(path)
    parts: List[str] = []

    # Read normal paragraphs
    for p in d.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)

    # Read table content
    for table in d.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    t = (p.text or "").strip()
                    if t:
                        parts.append(t)

    return "\n".join(parts)


def extract_pptx(path: str) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation(path)
    parts: List[str] = []

    for slide_number, slide in enumerate(prs.slides, 1):
        slide_parts: List[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    slide_parts.append(t)

        if slide_parts:
            parts.append(f"Slide {slide_number}:\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)


def extract_pdf(path: str) -> str:
    if PdfReader is None:
        raise RuntimeError("pypdf/PyPDF2 is not installed. Run: pip install pypdf")

    reader = PdfReader(path)
    parts: List[str] = []

    for page_number, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""

        text = text.strip()
        if text:
            parts.append(f"Page {page_number}:\n{text}")

    return "\n\n".join(parts)


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()

    if ext in {".txt", ".md"}:
        return read_text_file(path)

    if ext == ".docx":
        return extract_docx(path)

    if ext == ".pptx":
        return extract_pptx(path)

    if ext == ".pdf":
        return extract_pdf(path)

    return ""


# ---------------------------
# Tokenization / preprocessing
# ---------------------------

def tokenize(text: str) -> List[str]:
    tokens = TOKEN_RE.findall(text.lower())
    cleaned = [t for t in tokens if t not in STOPWORDS and len(t) >= 2]
    return cleaned


def split_sentences(text: str) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()

    if not text:
        return []

    sents = re.split(r"(?<=[\.\!\?])\s+", text)
    out: List[str] = []

    for s in sents:
        s = s.strip()
        if s:
            out.append(s)

    return out


def normalize_match_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).casefold().strip()


def extract_query_terms(query: str) -> List[str]:
    terms: List[str] = []

    for term in SEARCH_TERM_RE.findall(query.casefold()):
        term = term.strip()

        if not term:
            continue

        if term.isascii() and term in STOPWORDS:
            continue

        if term not in terms:
            terms.append(term)

    if not terms:
        normalized = normalize_match_text(query)
        if normalized:
            terms.append(normalized)

    return terms


def count_occurrences(text: str, needle: str) -> int:
    if not needle:
        return 0
    return text.count(needle)


def make_snippet(text: str, query_terms: List[str], max_len: int = 220) -> str:
    sents = split_sentences(text)

    if not sents:
        t = re.sub(r"\s+", " ", text).strip()
        return (t[:max_len] + "...") if len(t) > max_len else t

    qset = set(query_terms)
    best_sentence = ""
    best_score = 0

    for s in sents[:2000]:
        toks = set(tokenize(s))
        score = len(toks & qset)

        if score > best_score:
            best_sentence = s
            best_score = score

    snippet = best_sentence if best_sentence else sents[0]
    snippet = re.sub(r"\s+", " ", snippet).strip()

    if len(snippet) > max_len:
        snippet = snippet[:max_len] + "..."

    return snippet


def make_match_snippet(text: str, query: str, query_terms: List[str], max_len: int = 260) -> str:
    if not text:
        return ""

    lower_text = text.casefold()
    candidates = [query.casefold().strip(), *query_terms]
    candidates = sorted({c for c in candidates if c}, key=len, reverse=True)

    hit_at = -1

    for term in candidates:
        hit_at = lower_text.find(term)
        if hit_at >= 0:
            break

    if hit_at < 0:
        return make_snippet(text, query_terms, max_len=max_len)

    start = max(0, hit_at - max_len // 3)
    end = min(len(text), start + max_len)
    start = max(0, end - max_len)

    snippet = re.sub(r"\s+", " ", text[start:end]).strip()

    if start > 0:
        snippet = "..." + snippet

    if end < len(text):
        snippet += "..."

    return snippet


# ---------------------------
# TF-IDF building
# ---------------------------

def tf_counts(tokens: List[str]) -> Dict[str, int]:
    d: Dict[str, int] = {}

    for t in tokens:
        d[t] = d.get(t, 0) + 1

    return d


def compute_idf(docs_tokens: Dict[str, List[str]]) -> Dict[str, float]:
    n_docs = len(docs_tokens)
    df: Dict[str, int] = {}

    for toks in docs_tokens.values():
        seen = set(toks)

        for t in seen:
            df[t] = df.get(t, 0) + 1

    idf: Dict[str, float] = {}

    for term, dfi in df.items():
        idf[term] = math.log((n_docs + 1.0) / (dfi + 1.0)) + 1.0

    return idf


def build_vectors(
    docs_tokens: Dict[str, List[str]],
    idf: Dict[str, float],
) -> Tuple[Dict[str, Dict[str, float]], Dict[str, float]]:

    vectors: Dict[str, Dict[str, float]] = {}
    norms: Dict[str, float] = {}

    for doc_id, toks in docs_tokens.items():
        counts = tf_counts(toks)

        if not counts:
            vectors[doc_id] = {}
            norms[doc_id] = 0.0
            continue

        max_tf = max(counts.values())
        vec: Dict[str, float] = {}

        for term, tf in counts.items():
            weighted_tf = 0.5 + 0.5 * (tf / max_tf)
            weight = weighted_tf * idf.get(term, 0.0)

            if weight > 0:
                vec[term] = weight

        norm = math.sqrt(sum(v * v for v in vec.values()))

        vectors[doc_id] = vec
        norms[doc_id] = norm

    return vectors, norms


# ---------------------------
# Search
# ---------------------------

def cosine_sim(
    q: Dict[str, float],
    d: Dict[str, float],
    qn: float,
    dn: float,
) -> float:

    if qn <= 0 or dn <= 0:
        return 0.0

    if len(q) > len(d):
        q, d = d, q

    dot = 0.0

    for term, w in q.items():
        dot += w * d.get(term, 0.0)

    return dot / (qn * dn)


def build_query_vector(
    query: str,
    idf: Dict[str, float],
) -> Tuple[Dict[str, float], float, List[str]]:

    toks = tokenize(query)

    if not toks:
        return {}, 0.0, []

    counts = tf_counts(toks)
    max_tf = max(counts.values())

    qvec: Dict[str, float] = {}

    for term, tf in counts.items():
        weighted_tf = 0.5 + 0.5 * (tf / max_tf)
        qvec[term] = weighted_tf * idf.get(term, 0.0)

    qn = math.sqrt(sum(v * v for v in qvec.values()))

    return qvec, qn, toks


def search(index: TfIdfIndex, query: str, top_k: int = 5) -> List[Tuple[float, Document]]:
    qvec, qn, _qtoks = build_query_vector(query, index.idf)

    if not qvec:
        return []

    scored: List[Tuple[float, Document]] = []
    doc_map = {d.doc_id: d for d in index.documents}

    for doc_id, dvec in index.vectors.items():
        score = cosine_sim(qvec, dvec, qn, index.norms.get(doc_id, 0.0))

        if score > 0:
            scored.append((score, doc_map[doc_id]))

    scored.sort(key=lambda x: x[0], reverse=True)

    return scored[:top_k]


def text_match_search(index: TfIdfIndex, query: str, top_k: int = 50) -> List[SearchResult]:
    phrase = normalize_match_text(query)
    terms = extract_query_terms(query)

    if not phrase and not terms:
        return []

    results: List[SearchResult] = []

    for doc in index.documents:
        doc_text = normalize_match_text(doc.text)
        title_text = normalize_match_text(doc.title)

        phrase_hits = count_occurrences(doc_text, phrase) if phrase else 0
        term_hits = sum(count_occurrences(doc_text, term) for term in terms if term != phrase)

        title_hits = 0

        if phrase and phrase in title_text:
            title_hits += 1

        title_hits += sum(1 for term in terms if term != phrase and term in title_text)

        matches = phrase_hits + term_hits + title_hits

        if matches <= 0:
            continue

        score = phrase_hits * 20.0 + term_hits * 3.0 + title_hits * 10.0

        results.append(
            SearchResult(
                score=score,
                document=doc,
                snippet=make_match_snippet(doc.text, query, terms),
                matches=matches,
                source="Full-text match",
            )
        )

    results.sort(key=lambda r: r.score, reverse=True)

    return results[:top_k]


def smart_search(index: TfIdfIndex, query: str, top_k: int = 50) -> List[SearchResult]:
    merged: Dict[str, SearchResult] = {}

    for result in text_match_search(index, query, top_k=top_k * 2):
        merged[result.document.doc_id] = result

    tfidf_results = search(index, query, top_k=top_k * 2)
    _qvec, _qn, qtoks = build_query_vector(query, index.idf)

    for tfidf_score, doc in tfidf_results:
        if doc.doc_id in merged:
            existing = merged[doc.doc_id]
            existing.score += tfidf_score
            existing.source = "Full-text + keywords"
            continue

        merged[doc.doc_id] = SearchResult(
            score=tfidf_score,
            document=doc,
            snippet=make_snippet(doc.text, qtoks or extract_query_terms(query)),
            matches=0,
            source="Keyword relevance",
        )

    results = list(merged.values())
    results.sort(key=lambda r: r.score, reverse=True)

    return results[:top_k]


# ---------------------------
# Indexing / IO
# ---------------------------

def walk_files(root: str) -> List[str]:
    paths: List[str] = []

    for base, _dirs, files in os.walk(root):
        for filename in files:
            ext = os.path.splitext(filename)[1].lower()

            if ext in SUPPORTED_EXT:
                paths.append(os.path.join(base, filename))

    return paths


def doc_title_from_path(path: str) -> str:
    return os.path.splitext(os.path.basename(path))[0]


def build_index(
    root_folder: str,
    verbose: bool = True,
    progress_callback: Optional[Callable[[int, int, str], None]] = None,
) -> TfIdfIndex:

    root_folder = os.path.abspath(root_folder)
    paths = walk_files(root_folder)

    documents: List[Document] = []
    docs_tokens: Dict[str, List[str]] = {}

    if verbose:
        print(f"Indexing folder: {root_folder}")
        print(f"Found {len(paths)} supported files.")

    if progress_callback is not None:
        progress_callback(0, len(paths), "Preparing files")

    for i, path in enumerate(paths, 1):
        if progress_callback is not None and (i == 1 or i == len(paths) or i % 5 == 0):
            progress_callback(i, len(paths), os.path.basename(path))

        title = doc_title_from_path(path)

        try:
            text = extract_text(path)
        except Exception as e:
            if verbose:
                print(f"[WARN] Cannot read {path}: {e}")
            text = ""

        text = text.strip()

        # Important fallback:
        # If the file body cannot be extracted, still index the filename.
        if not text:
            text = title

        doc_id = f"d{i}_{abs(hash(path)) % 10_000_000}"
        toks = tokenize(text + " " + title)

        doc = Document(
            doc_id=doc_id,
            path=path,
            title=title,
            text=text,
            n_chars=len(text),
            n_tokens=len(toks),
        )

        documents.append(doc)
        docs_tokens[doc_id] = toks

    if progress_callback is not None:
        progress_callback(len(paths), len(paths), "Calculating keyword weights")

    idf = compute_idf(docs_tokens)
    vectors, norms = build_vectors(docs_tokens, idf)

    meta = IndexMeta(
        created_at=time.time(),
        root_folder=root_folder,
        n_docs=len(documents),
    )

    return TfIdfIndex(
        meta=meta,
        documents=documents,
        idf=idf,
        vectors=vectors,
        norms=norms,
    )


def save_index(index: TfIdfIndex, path: str) -> None:
    data = {
        "meta": asdict(index.meta),
        "documents": [asdict(d) for d in index.documents],
        "idf": index.idf,
        "vectors": index.vectors,
        "norms": index.norms,
    }

    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"Saved index to: {path}")


def load_index(path: str) -> TfIdfIndex:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)

    meta = IndexMeta(**data["meta"])
    documents = [Document(**d) for d in data["documents"]]

    return TfIdfIndex(
        meta=meta,
        documents=documents,
        idf=data["idf"],
        vectors=data["vectors"],
        norms=data["norms"],
    )


# ---------------------------
# Desktop GUI
# ---------------------------

class StudySearchApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Local Study Material Search")
        self.root.geometry("1100x720")
        self.root.minsize(860, 560)

        default_folder = os.path.join(os.path.expanduser("~"), "Documents")

        if not os.path.isdir(default_folder):
            default_folder = os.getcwd()

        self.folder_var = tk.StringVar(value=default_folder)
        self.query_var = tk.StringVar()
        self.status_var = tk.StringVar(value="Please choose a materials folder")
        self.progress_var = tk.DoubleVar(value=0.0)

        self.index: Optional[TfIdfIndex] = None
        self.index_folder = ""
        self.results: List[SearchResult] = []
        self.last_query = ""

        self.worker: Optional[threading.Thread] = None
        self.search_worker: Optional[threading.Thread] = None
        self.busy_buttons: List[object] = []

        self.build_widgets()

    def build_widgets(self) -> None:
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)

        main = ttk.Frame(self.root, padding=12)
        main.grid(row=0, column=0, sticky="nsew")
        main.columnconfigure(1, weight=1)
        main.rowconfigure(2, weight=1)

        ttk.Label(main, text="Materials Folder").grid(
            row=0,
            column=0,
            sticky="w",
            padx=(0, 8),
            pady=(0, 8),
        )

        folder_entry = ttk.Entry(main, textvariable=self.folder_var)
        folder_entry.grid(row=0, column=1, sticky="ew", pady=(0, 8))

        browse_btn = ttk.Button(main, text="Browse", command=self.browse_folder)
        browse_btn.grid(row=0, column=2, sticky="ew", padx=(8, 0), pady=(0, 8))

        refresh_btn = ttk.Button(main, text="Refresh Index", command=self.refresh_index)
        refresh_btn.grid(row=0, column=3, sticky="ew", padx=(8, 0), pady=(0, 8))

        ttk.Label(main, text="Search").grid(
            row=1,
            column=0,
            sticky="w",
            padx=(0, 8),
            pady=(0, 8),
        )

        query_entry = ttk.Entry(main, textvariable=self.query_var)
        query_entry.grid(row=1, column=1, sticky="ew", pady=(0, 8))
        query_entry.bind("<Return>", lambda _event: self.search_requested())

        search_btn = ttk.Button(main, text="Search", command=self.search_requested)
        search_btn.grid(row=1, column=2, columnspan=2, sticky="ew", padx=(8, 0), pady=(0, 8))

        columns = ("rank", "score", "title", "matches", "source", "path")
        self.tree = ttk.Treeview(main, columns=columns, show="headings", selectmode="browse")

        self.tree.heading("rank", text="Rank")
        self.tree.heading("score", text="Score")
        self.tree.heading("title", text="File Name")
        self.tree.heading("matches", text="Matches")
        self.tree.heading("source", text="Type")
        self.tree.heading("path", text="Path")

        self.tree.column("rank", width=60, anchor="center", stretch=False)
        self.tree.column("score", width=80, anchor="center", stretch=False)
        self.tree.column("title", width=220, stretch=True)
        self.tree.column("matches", width=80, anchor="center", stretch=False)
        self.tree.column("source", width=140, anchor="center", stretch=False)
        self.tree.column("path", width=520, stretch=True)

        self.tree.grid(row=2, column=0, columnspan=4, sticky="nsew")
        self.tree.bind("<<TreeviewSelect>>", lambda _event: self.show_selected_detail())
        self.tree.bind("<Double-1>", lambda _event: self.open_selected_file())

        scroll = ttk.Scrollbar(main, orient="vertical", command=self.tree.yview)
        self.tree.configure(yscrollcommand=scroll.set)
        scroll.grid(row=2, column=4, sticky="ns")

        actions = ttk.Frame(main)
        actions.grid(row=3, column=0, columnspan=4, sticky="ew", pady=(8, 8))

        open_btn = ttk.Button(actions, text="Open File", command=self.open_selected_file)
        open_btn.pack(side="left")

        copy_btn = ttk.Button(actions, text="Copy Path", command=self.copy_selected_path)
        copy_btn.pack(side="left", padx=(8, 0))

        export_btn = ttk.Button(actions, text="Export Results", command=self.export_results)
        export_btn.pack(side="left", padx=(8, 0))

        self.detail = tk.Text(main, height=8, wrap="word")
        self.detail.grid(row=4, column=0, columnspan=4, sticky="nsew")
        self.detail.configure(state="disabled")

        self.progress = ttk.Progressbar(
            main,
            variable=self.progress_var,
            maximum=100,
            mode="determinate",
        )
        self.progress.grid(row=5, column=0, columnspan=4, sticky="ew", pady=(8, 0))

        status = ttk.Label(main, textvariable=self.status_var, anchor="w")
        status.grid(row=6, column=0, columnspan=4, sticky="ew", pady=(8, 0))

        self.busy_buttons = [
            browse_btn,
            refresh_btn,
            search_btn,
            open_btn,
            copy_btn,
            export_btn,
        ]

        query_entry.focus_set()

    def set_busy(self, busy: bool) -> None:
        state = "disabled" if busy else "normal"

        for button in self.busy_buttons:
            button.configure(state=state)

        self.root.configure(cursor="watch" if busy else "")
        self.root.update_idletasks()

    def browse_folder(self) -> None:
        folder = filedialog.askdirectory(initialdir=self.folder_var.get() or os.getcwd())

        if not folder:
            return

        self.folder_var.set(folder)

        if os.path.abspath(folder) != self.index_folder:
            self.index = None
            self.index_folder = ""
            self.progress_var.set(0)
            self.status_var.set("Folder selected. Enter keywords and search.")

    def refresh_index(self) -> None:
        self.start_indexing()

    def update_index_progress(self, current: int, total: int, filename: str) -> None:
        if total <= 0:
            self.progress_var.set(0)
            self.status_var.set("No supported files found")
            return

        percent = max(0.0, min(100.0, current * 100.0 / total))
        self.progress_var.set(percent)

        if current >= total:
            self.status_var.set("Calculating search index...")
        else:
            self.status_var.set(f"Building index: {current}/{total}  {filename}")

    def search_requested(self) -> None:
        query = self.query_var.get().strip()

        if not query:
            messagebox.showinfo("Search", "Please enter a keyword or sentence.")
            return

        if self.search_worker is not None and self.search_worker.is_alive():
            self.status_var.set("Searching. Please wait.")
            return

        folder = os.path.abspath(self.folder_var.get().strip())

        if self.index is None or self.index_folder != folder:
            self.start_indexing(after_query=query)
            return

        self.start_search(query)

    def start_indexing(self, after_query: str = "") -> None:
        folder = os.path.abspath(self.folder_var.get().strip())

        if not os.path.isdir(folder):
            messagebox.showerror("Materials Folder", "Please choose an existing folder.")
            return

        if self.worker is not None and self.worker.is_alive():
            self.status_var.set("Index is being built. Please wait.")
            return

        self.set_busy(True)
        self.progress_var.set(0)
        self.status_var.set("Reading files and building index...")

        def run() -> None:
            def progress(current: int, total: int, filename: str) -> None:
                self.root.after(
                    0,
                    lambda: self.update_index_progress(current, total, filename),
                )

            try:
                index = build_index(folder, verbose=False, progress_callback=progress)
                error = None
            except Exception as exc:
                index = None
                error = exc

            self.root.after(
                0,
                lambda: self.finish_indexing(folder, index, error, after_query),
            )

        self.worker = threading.Thread(target=run, daemon=True)
        self.worker.start()

    def finish_indexing(
        self,
        folder: str,
        index: Optional[TfIdfIndex],
        error: Optional[Exception],
        after_query: str,
    ) -> None:

        self.set_busy(False)

        if error is not None or index is None:
            messagebox.showerror("Index Failed", str(error))
            self.status_var.set("Index failed")
            return

        self.index = index
        self.index_folder = folder
        self.progress_var.set(100)

        self.status_var.set(f"Index completed: {index.meta.n_docs} files")

        if index.meta.n_docs == 0:
            messagebox.showwarning(
                "Empty Index",
                "No files were successfully indexed.\n\n"
                "Please check:\n"
                "1. The folder contains .docx / .pptx / .pdf / .txt / .md files\n"
                "2. Required packages are installed: python-docx, python-pptx, pypdf\n"
                "3. The PDF is not only a scanned image\n"
                "4. The files are not old .doc or .ppt formats",
            )

        if after_query:
            self.start_search(after_query)

    def start_search(self, query: str) -> None:
        if self.index is None:
            return

        if self.search_worker is not None and self.search_worker.is_alive():
            self.status_var.set("Searching. Please wait.")
            return

        self.last_query = query
        current_index = self.index

        self.set_busy(True)
        self.status_var.set("Searching document content...")
        self.progress.configure(mode="indeterminate")
        self.progress.start(12)

        def run() -> None:
            try:
                results = smart_search(current_index, query, top_k=100)
                error = None
            except Exception as exc:
                results = []
                error = exc

            self.root.after(0, lambda: self.finish_search(query, results, error))

        self.search_worker = threading.Thread(target=run, daemon=True)
        self.search_worker.start()

    def finish_search(
        self,
        query: str,
        results: List[SearchResult],
        error: Optional[Exception],
    ) -> None:

        self.progress.stop()
        self.progress.configure(mode="determinate")
        self.progress_var.set(100)
        self.set_busy(False)

        if error is not None:
            messagebox.showerror("Search Failed", str(error))
            self.status_var.set("Search failed")
            return

        self.render_search_results(query, results)

    def render_search_results(self, query: str, results: List[SearchResult]) -> None:
        self.last_query = query
        self.results = results

        for item in self.tree.get_children():
            self.tree.delete(item)

        for rank, result in enumerate(self.results, 1):
            self.tree.insert(
                "",
                "end",
                iid=str(rank - 1),
                values=(
                    rank,
                    f"{result.score:.2f}",
                    result.document.title,
                    result.matches,
                    result.source,
                    result.document.path,
                ),
            )

        if self.results:
            first = self.tree.get_children()[0]
            self.tree.selection_set(first)
            self.tree.focus(first)
            self.show_selected_detail()
            self.status_var.set(f"Found {len(self.results)} results")
        else:
            self.clear_detail()
            self.status_var.set("No matching results found")

    def selected_result(self) -> Optional[SearchResult]:
        selected = self.tree.selection()

        if not selected:
            return None

        try:
            return self.results[int(selected[0])]
        except Exception:
            return None

    def clear_detail(self) -> None:
        self.detail.configure(state="normal")
        self.detail.delete("1.0", tk.END)
        self.detail.configure(state="disabled")

    def show_selected_detail(self) -> None:
        result = self.selected_result()

        if result is None:
            self.clear_detail()
            return

        text = (
            f"File Name: {result.document.title}\n"
            f"Path: {result.document.path}\n"
            f"Matches: {result.matches}    "
            f"Type: {result.source}    "
            f"Score: {result.score:.2f}\n\n"
            f"{result.snippet}"
        )

        self.detail.configure(state="normal")
        self.detail.delete("1.0", tk.END)
        self.detail.insert("1.0", text)
        self.detail.configure(state="disabled")

    def open_selected_file(self) -> None:
        result = self.selected_result()

        if result is None:
            return

        path = result.document.path

        if not os.path.exists(path):
            messagebox.showerror("Open File", "File does not exist.")
            return

        try:
            if sys.platform.startswith("win"):
                os.startfile(path)  # type: ignore[attr-defined]
            elif sys.platform == "darwin":
                subprocess.Popen(["open", path])
            else:
                subprocess.Popen(["xdg-open", path])
        except Exception as exc:
            messagebox.showerror("Open File", str(exc))

    def copy_selected_path(self) -> None:
        result = self.selected_result()

        if result is None:
            return

        self.root.clipboard_clear()
        self.root.clipboard_append(result.document.path)
        self.status_var.set("File path copied")

    def export_results(self) -> None:
        if not self.results:
            messagebox.showinfo("Export Results", "There are no search results to export.")
            return

        out_path = filedialog.asksaveasfilename(
            title="Export Search Results",
            defaultextension=".md",
            filetypes=[
                ("Markdown", "*.md"),
                ("Text", "*.txt"),
                ("All files", "*.*"),
            ],
        )

        if not out_path:
            return

        lines = [
            "# Search Results",
            f"- Query: **{self.last_query}**",
            "",
        ]

        for i, result in enumerate(self.results, 1):
            lines.append(f"## {i}. {result.document.title}")
            lines.append(f"- Score: `{result.score:.4f}`")
            lines.append(f"- Matches: `{result.matches}`")
            lines.append(f"- Type: `{result.source}`")
            lines.append(f"- Path: `{result.document.path}`")
            lines.append(f"- Snippet: {result.snippet}")
            lines.append("")

        with open(out_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))

        self.status_var.set(f"Exported to: {out_path}")


def launch_gui() -> int:
    if tk is None or ttk is None:
        print("Tkinter is not available in this Python environment.")
        return 1

    root = tk.Tk()
    StudySearchApp(root)
    root.mainloop()

    return 0


# ---------------------------
# Interactive CLI
# ---------------------------

HELP_TEXT = """
Commands:
  help                         Show help
  info                         Show index info
  query <text> [-k N]          Search for text
  open <rank>                  Print full text of a result from last query
  snippet <rank>               Print snippet of a result from last query
  export <file.md>             Export last query results to Markdown
  exit                         Quit

Examples:
  query medical chatbot -k 5
  open 1
  export results.md
"""


def print_info(index: TfIdfIndex) -> None:
    created = datetime.fromtimestamp(index.meta.created_at).strftime("%Y-%m-%d %H:%M")

    print("Index Info")
    print("-" * 60)
    print(f"Root folder: {index.meta.root_folder}")
    print(f"Created at : {created}")
    print(f"Documents  : {index.meta.n_docs}")

    rare = sorted(index.idf.items(), key=lambda x: x[1], reverse=True)[:10]

    print("Top rare terms:", ", ".join(t for t, _ in rare))


def export_results_md(
    out_path: str,
    query: str,
    results: List[SearchResult],
) -> None:

    lines: List[str] = []
    lines.append("# Search Results")
    lines.append(f"- Query: **{query}**")
    lines.append("")

    if not results:
        lines.append("_No matches found._")
    else:
        for i, result in enumerate(results, 1):
            lines.append(f"## {i}. {result.document.title}")
            lines.append(f"- Score: `{result.score:.4f}`")
            lines.append(f"- Matches: `{result.matches}`")
            lines.append(f"- Type: `{result.source}`")
            lines.append(f"- Path: `{result.document.path}`")
            lines.append(f"- Snippet: {result.snippet}")
            lines.append("")

    with open(out_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"Exported to: {out_path}")


def parse_k(args: List[str], default_k: int = 5) -> int:
    if "-k" in args:
        try:
            idx = args.index("-k")
            return max(1, int(args[idx + 1]))
        except Exception:
            return default_k

    return default_k


def remove_k_args(args: List[str]) -> List[str]:
    q_tokens = []
    skip = 0

    for i, a in enumerate(args):
        if skip:
            skip -= 1
            continue

        if a == "-k" and i + 1 < len(args):
            skip = 1
            continue

        q_tokens.append(a)

    return q_tokens


def repl(index: TfIdfIndex) -> None:
    print("Local Study Search Engine")
    print("Type 'help' for commands.")
    print("-" * 60)

    last_query = ""
    last_results: List[SearchResult] = []

    while True:
        try:
            raw = input("> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye.")
            break

        if not raw:
            continue

        parts = raw.split()
        cmd = parts[0].lower()
        args = parts[1:]

        if cmd in ("exit", "quit"):
            print("Bye.")
            break

        if cmd == "help":
            print(HELP_TEXT)
            continue

        if cmd == "info":
            print_info(index)
            continue

        if cmd == "query":
            if not args:
                print("Please provide query text.")
                continue

            k = parse_k(args, default_k=5)
            query_text = " ".join(remove_k_args(args)).strip()

            last_query = query_text
            last_results = smart_search(index, query_text, top_k=k)

            if not last_results:
                print("No matches found.")
                continue

            print(f"Results for: {query_text}")
            print("-" * 60)

            for i, result in enumerate(last_results, 1):
                print(f"{i}. {result.document.title} | score={result.score:.4f}")
                print(f"   Type: {result.source} | Matches: {result.matches}")
                print(f"   Path: {result.document.path}")
                print(f"   {result.snippet}")

            continue

        if cmd == "open":
            if not last_results:
                print("Run a query first.")
                continue

            if not args:
                print("Usage: open <rank>")
                continue

            try:
                r = max(1, int(args[0]))
            except Exception:
                print("Rank must be a number.")
                continue

            if r > len(last_results):
                print("Rank out of range.")
                continue

            doc = last_results[r - 1].document

            print("-" * 60)
            print(f"OPEN: {doc.title}")
            print(f"PATH: {doc.path}")
            print("-" * 60)
            print(doc.text[:8000])

            if len(doc.text) > 8000:
                print("\n[... truncated ...]")

            continue

        if cmd == "snippet":
            if not last_results:
                print("Run a query first.")
                continue

            if not args:
                print("Usage: snippet <rank>")
                continue

            try:
                r = max(1, int(args[0]))
            except Exception:
                print("Rank must be a number.")
                continue

            if r > len(last_results):
                print("Rank out of range.")
                continue

            print(last_results[r - 1].snippet)
            continue

        if cmd == "export":
            if not last_results:
                print("Run a query first.")
                continue

            out = args[0] if args else "results.md"
            export_results_md(out, last_query, last_results)
            continue

        print("Unknown command. Type 'help'.")


# ---------------------------
# Main
# ---------------------------

def usage() -> None:
    print(__doc__)


def main(argv: List[str]) -> int:
    if len(argv) < 2:
        return launch_gui()

    mode = argv[1].lower()

    if mode == "gui":
        return launch_gui()

    if mode == "index":
        if len(argv) < 3:
            print("Usage: python study_search.py index <folder_path> [index.json]")
            return 1

        folder = argv[2]
        out = argv[3] if len(argv) >= 4 else "index.json"

        idx = build_index(folder)
        save_index(idx, out)

        return 0

    if mode == "search":
        idx_path = argv[2] if len(argv) >= 3 else "index.json"
        idx = load_index(idx_path)
        repl(idx)

        return 0

    if mode == "query":
        if len(argv) < 4:
            print('Usage: python study_search.py query <index.json> "<query>" [-k 5]')
            return 1

        idx_path = argv[2]
        query_text = argv[3]
        k = 5

        if "-k" in argv:
            try:
                k = int(argv[argv.index("-k") + 1])
            except Exception:
                k = 5

        idx = load_index(idx_path)
        results = smart_search(idx, query_text, top_k=k)

        if not results:
            print("No matches found.")
            return 0

        print(f"Results for: {query_text}")
        print("-" * 60)

        for i, result in enumerate(results, 1):
            print(f"{i}. {result.document.title} | score={result.score:.4f}")
            print(f"   Type: {result.source} | Matches: {result.matches}")
            print(f"   Path: {result.document.path}")
            print(f"   {result.snippet}")

        return 0

    usage()
    return 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))